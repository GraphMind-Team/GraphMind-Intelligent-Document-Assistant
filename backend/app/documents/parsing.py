"""PDF/Markdown/HTML parsing and chunking (Story 2.3, FR-3).

Output is plain text only, all the way through -- extracted chapter
titles and body text are never rendered as HTML or Markdown anywhere
downstream (a standing constraint: a malicious `.md`/`.html` upload's
heading text must never reach a raw-HTML renderer, per
`_bmad-output/implementation-artifacts/deferred-work.md`).

Uploaded bytes are only extension/content-type-checked at upload time,
never magic-byte-verified (same doc) -- every parser here must catch its
own failure mode and raise `UnparseableDocument` rather than assume the
bytes are trustworthy for their claimed `file_type`.
"""

import copy
import io
import re
from dataclasses import dataclass
from typing import Final
from zipfile import BadZipFile, ZipFile

from bs4 import BeautifulSoup, NavigableString
from bs4.element import PreformattedString
from docx import Document as DocxDocument
from docx.table import Table as DocxTable
from pptx import Presentation
from pypdf import PdfReader

# ~250 words per chunk with a ~40-word overlap between consecutive chunks
# -- cheap, and avoids losing a passage that straddles a chunk boundary.
# chunk_index is sequential across the whole document, not reset per
# chapter.
#
# Sized against the embedding model's input limit, not picked
# independently of it: English text tokenizes at roughly 1.3 tokens/word,
# but the model is multilingual, and Cyrillic (Bulgarian) text can
# tokenize considerably denser under a shared multilingual vocabulary --
# 250 words leaves real headroom even at a pessimistic ~2 tokens/word,
# rather than assuming the English ratio holds for every document.
# Getting this wrong is silent, not an error: the model truncates instead
# of failing, so an oversized chunk just means the tail of its (stored,
# citable) text was never actually seen by its own embedding.
#
# The limit this was originally sized against was 512 tokens, from the
# in-process fastembed model. Embedding moved to Weaviate
# (text2vec-weaviate / arctic-embed-l-v2.0, see
# `weaviate_client.EMBEDDING_MODEL`), whose context window is far larger,
# so 250 words is now comfortably inside it rather than near the edge.
# Left unchanged deliberately: `CHUNK_OVERLAP_WORDS` below, citation
# granularity, and Story 2.4's extraction concatenation are all tuned
# around this size, so growing it is a retrieval-quality decision to make
# on its own evidence, not a free consequence of the bigger window.
_CHUNK_WORD_COUNT = 250

# Public (not `_`-prefixed), for the same reason `weaviate_client
# .PASSAGE_BATCH_SIZE` is: `service.py` strips this overlap back off when
# it concatenates chunks for entity extraction (Story 2.4), so the two
# sides read one source of truth instead of two magic numbers that could
# silently drift apart.
CHUNK_OVERLAP_WORDS = 40

_FULL_DOCUMENT_CHAPTER = "Full Document"
_HEADING_TAGS = ("h1", "h2", "h3")

_MARKDOWN_HEADING_RE = re.compile(r"^#{1,6}\s+(.*)$", re.MULTILINE)

# Fenced code blocks (```...```) can contain lines that look like ATX
# headings (e.g. a `# install deps` comment in a bash snippet) but aren't
# -- matched separately so heading detection can skip anything inside one.
# Non-greedy DOTALL match from an opening ``` line to the next ``` line.
_MARKDOWN_CODE_FENCE_RE = re.compile(r"^```.*?^```[ \t]*$", re.MULTILINE | re.DOTALL)

# DOCX/PPTX are zip archives -- unlike every other supported format, the
# 20MB `MAX_FILE_SIZE_BYTES` upload check only bounds the *compressed*
# size. DEFLATE lets a small archive expand enormously (a few hundred KB
# of crafted XML can unpack to tens of GB), and python-docx/python-pptx
# both fully materialize the inflated XML tree in memory before either
# library gets a chance to fail gracefully -- so this has to be checked
# before either library opens the file at all. 200MB is generous headroom
# above what a legitimate multi-hundred-page report or image-heavy deck
# expands to (embedded images are already compressed inside the zip,
# close to 1:1), while still bounding a worker's memory use well below a
# crafted bomb's thousand-to-one ratio.
_MAX_ZIP_UNCOMPRESSED_BYTES: Final = 200 * 1024 * 1024

# Read (and immediately discard) each entry's decompressed output in
# fixed-size chunks -- never the whole entry at once -- so the check
# itself never holds more than one chunk of inflated data in memory
# regardless of how large the real payload turns out to be.
_ZIP_READ_CHUNK_SIZE: Final = 1024 * 1024

# Stand-in `ZipInfo.file_size` used only while measuring an entry (see
# `_check_zip_bomb`), large enough that `zipfile`'s per-read truncation
# against it can never fire -- the entry's real deflate stream ends long
# before this, whatever it declared about itself. Not `sys.maxsize`:
# `zipfile` writes this field back out as a 32/64-bit integer in some
# paths, so an absurd-but-representable value is the safer stand-in.
_UNTRUNCATED_ZIP_ENTRY_SIZE: Final = 1 << 62


def _check_zip_bomb(content: bytes) -> None:
    """Raises `UnparseableDocument` if `content` (a DOCX/PPTX zip archive)
    decompresses to more than `_MAX_ZIP_UNCOMPRESSED_BYTES`.

    Deliberately does *not* trust `ZipInfo.file_size` as a stand-in for
    this -- that field is read straight from the archive's central
    directory, which is ordinary attacker-controlled archive data, not a
    verified measurement of anything. So this runs every entry through
    the real decompressor instead, in fixed-size chunks, and counts what
    actually comes out; the check is only ever as trustworthy as the
    bytes it has itself inflated and counted.

    Streaming the entry is not by itself enough to stop trusting the
    declared size, because `zipfile` applies that size to the *read* as
    well: `ZipExtFile` truncates each read to the declared byte count it
    has left (`data = data[:self._left]`). Reading an entry normally
    would therefore measure the lie rather than the payload -- an entry
    declaring 1KB behind a 400MB deflate stream gets counted as 1KB and
    waved through. Crucially that truncation is applied *after* the
    inflate, not instead of it, so it bounds what a reader receives but
    not what decompressing allocated to produce it: `ZipFile.read(name)`
    (exactly what python-docx/python-pptx call) passes no size, which
    becomes `max_length = 1 << 31` into the decompressor -- effectively
    unbounded -- so the full 400MB is materialized and only then sliced
    back down to the declared 1KB. Measured end to end, a 433KB upload
    crafted this way peaked at 830MB of heap with the declared size
    trusted. Forging the entry's CRC-32 to match the truncated prefix
    isn't even required for that; without it the `BadZipFile` is raised
    *after* the allocation, so it changes the error message and nothing
    about the memory.

    Hence `probe.file_size` below: measuring against a copy of the
    `ZipInfo` whose declared size can't truncate anything is what makes
    the count reflect the real deflate output. An honest entry is
    unaffected -- its stream simply ends where it always did, and its
    CRC-32 still checks out -- while an entry lying downward is now
    inflated under this function's own bounded chunk reads and rejected
    on what actually comes out of it.
    """
    try:
        with ZipFile(io.BytesIO(content)) as archive:
            total_uncompressed = 0
            for info in archive.infolist():
                # A copy, not `info` itself: `archive.infolist()` hands
                # out the live `ZipInfo` objects the `ZipFile` keeps for
                # its own reads, so mutating one here would corrupt the
                # archive's view of that entry for every later reader.
                probe = copy.copy(info)
                probe.file_size = _UNTRUNCATED_ZIP_ENTRY_SIZE
                with archive.open(probe) as entry:
                    while True:
                        chunk = entry.read(_ZIP_READ_CHUNK_SIZE)
                        if not chunk:
                            break
                        total_uncompressed += len(chunk)
                        if total_uncompressed > _MAX_ZIP_UNCOMPRESSED_BYTES:
                            raise UnparseableDocument(
                                "Document expands to an unreasonable size when decompressed."
                            )
    except BadZipFile as exc:
        raise UnparseableDocument("Could not parse archive.") from exc


@dataclass(frozen=True)
class ParsedChunk:
    chapter: str
    chunk_index: int
    text: str
    # True only for the first chunk produced within one chapter -- i.e. the
    # one chunk in this chapter that carries no `CHUNK_OVERLAP_WORDS`
    # carried over from a previous chunk. `service.py`'s
    # `_build_extraction_text` uses this (not a `chapter == previous_chapter`
    # title comparison) to decide whether to strip overlap back off before
    # sending text to the LLM: two distinct chapters can share the exact
    # same title (e.g. two "Overview" sections under different parts), which
    # a title comparison can't tell apart from one chapter continuing --
    # this flag is set here, by the only code that actually knows whether
    # overlap was applied, so it can't be fooled by that.
    is_chapter_start: bool = True


class UnparseableDocument(Exception):
    """`content` couldn't be parsed as its claimed `file_type`, or yielded
    no extractable text. The expected outcome for a corrupt or mislabeled
    upload -- callers must catch this, not let it propagate."""


def parse_document(file_type: str, content: bytes) -> list[ParsedChunk]:
    """PDF/Markdown/HTML bytes -> chaptered, chunked plain-text passages.

    Raises `UnparseableDocument` if the bytes can't be parsed as
    `file_type`, or if parsing succeeds but yields zero extractable text
    (an empty file or a scanned/image-only PDF) -- AC1 requires "one or
    more passages", so a silent zero-passage success is treated as a
    failure, not success.
    """
    if file_type == "pdf":
        chapters = _parse_pdf(content)
    elif file_type == "markdown":
        chapters = _parse_markdown(content)
    elif file_type == "html":
        chapters = _parse_html(content)
    elif file_type == "docx":
        chapters = _parse_docx(content)
    elif file_type == "pptx":
        chapters = _parse_pptx(content)
    else:
        raise UnparseableDocument(f"Unsupported file_type: {file_type!r}")

    chunks = _chunk_chapters(chapters)
    if not chunks:
        raise UnparseableDocument("Document produced no extractable text.")
    return chunks


def _chunk_chapters(chapters: list[tuple[str, str]]) -> list[ParsedChunk]:
    chunks: list[ParsedChunk] = []
    chunk_index = 0
    for chapter, text in chapters:
        words = text.split()
        if not words:
            # A heading with nothing under it -- a title-only slide, a
            # section-divider slide, a DOCX heading immediately followed
            # by another heading -- would otherwise vanish from the index
            # entirely rather than degrade to "no body text under this
            # title". The generic "Full Document" placeholder is the one
            # exception: that's not a real heading a user wrote, so an
            # empty one means the document (or this section of it) truly
            # had no extractable text, not a title worth indexing on its
            # own.
            if chapter == _FULL_DOCUMENT_CHAPTER:
                continue
            words = chapter.split()
            if not words:
                continue
        start = 0
        is_first_chunk_in_chapter = True
        while start < len(words):
            end = start + _CHUNK_WORD_COUNT
            chunk_text = " ".join(words[start:end])
            chunks.append(
                ParsedChunk(
                    chapter=chapter,
                    chunk_index=chunk_index,
                    text=chunk_text,
                    is_chapter_start=is_first_chunk_in_chapter,
                )
            )
            chunk_index += 1
            is_first_chunk_in_chapter = False
            if end >= len(words):
                break
            start = end - CHUNK_OVERLAP_WORDS
    return chunks


# --- Markdown ---------------------------------------------------------


def _parse_markdown(content: bytes) -> list[tuple[str, str]]:
    try:
        text = content.decode("utf-8")
    except UnicodeDecodeError as exc:
        raise UnparseableDocument("Could not decode Markdown as UTF-8.") from exc

    fenced_spans = [(m.start(), m.end()) for m in _MARKDOWN_CODE_FENCE_RE.finditer(text)]

    def _inside_a_fence(position: int) -> bool:
        return any(start <= position < end for start, end in fenced_spans)

    # A `# heading`-shaped line inside a fenced code block (e.g. a `#
    # install deps` comment in a bash snippet) is code, not a real
    # section boundary -- README-style uploads with shell examples hit
    # this often enough that it's worth filtering rather than ignoring.
    matches = [
        m for m in _MARKDOWN_HEADING_RE.finditer(text) if not _inside_a_fence(m.start())
    ]
    if not matches:
        return [(_FULL_DOCUMENT_CHAPTER, text)]

    chapters: list[tuple[str, str]] = []
    preamble = text[: matches[0].start()].strip()
    if preamble:
        # Text before the first heading (an abstract/intro paragraph, a
        # README's lead-in) is real content -- often exactly what answers
        # "what is this document about" -- not something to drop silently.
        chapters.append((_FULL_DOCUMENT_CHAPTER, preamble))
    for i, match in enumerate(matches):
        title = match.group(1).strip() or _FULL_DOCUMENT_CHAPTER
        body_start = match.end()
        body_end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        chapters.append((title, text[body_start:body_end].strip()))
    return chapters


# --- HTML ---------------------------------------------------------------


def _parse_html(content: bytes) -> list[tuple[str, str]]:
    try:
        text = content.decode("utf-8")
    except UnicodeDecodeError as exc:
        raise UnparseableDocument("Could not decode HTML as UTF-8.") from exc

    try:
        soup = BeautifulSoup(text, "html.parser")
    except Exception as exc:  # pragma: no cover - html.parser rarely raises
        raise UnparseableDocument("Could not parse HTML.") from exc

    for tag in soup(["script", "style"]):
        tag.decompose()

    root = soup.body or soup

    chapters: list[tuple[str, list[str]]] = []
    current_title = _FULL_DOCUMENT_CHAPTER
    current_parts: list[str] = []

    for element in root.descendants:
        if getattr(element, "name", None) in _HEADING_TAGS:
            # Same "keep a heading with no body under it" reasoning as
            # `_parse_docx` -- a heading directly followed by another
            # heading would otherwise disappear entirely instead of
            # falling through to `_chunk_chapters`'s empty-body handling.
            if current_parts or current_title != _FULL_DOCUMENT_CHAPTER:
                chapters.append((current_title, current_parts))
            current_title = element.get_text(strip=True) or _FULL_DOCUMENT_CHAPTER
            current_parts = []
        elif isinstance(element, NavigableString):
            # Comment/CData/Doctype/Declaration/ProcessingInstruction all
            # subclass NavigableString (via PreformattedString) in bs4 --
            # excluded here so an HTML comment (internal notes, TODOs,
            # tool leftovers -- never meant as visible content) doesn't
            # become embedded, indexed, extractable text.
            if isinstance(element, PreformattedString):
                continue
            # Skip the heading's own text node -- already captured as
            # current_title above, would otherwise duplicate into the body.
            if getattr(element.parent, "name", None) in _HEADING_TAGS:
                continue
            piece = str(element).strip()
            if piece:
                current_parts.append(piece)

    if current_parts or current_title != _FULL_DOCUMENT_CHAPTER:
        chapters.append((current_title, current_parts))

    return [(title, " ".join(parts)) for title, parts in chapters]


# --- DOCX -----------------------------------------------------------------

# Word's built-in heading styles, matched two ways:
#
# `style_id` is the locale-invariant identifier Word itself assigns on the
# style definition (`w:styleId`) -- stable regardless of the Word install's
# UI language, because it's fixed by the OOXML built-in style table rather
# than translated for display.
_DOCX_HEADING_STYLE_IDS = {"Title", "Heading1", "Heading2", "Heading3"}

# `style.name` is the human-readable label (`w:name`), which a
# non-English-localized Word *does* translate (e.g. German "Uberschrift 1"
# instead of "Heading 1") -- kept as a second check for documents from
# tools that set a recognizable name without the matching built-in
# `style_id` (exports from non-Word editors, hand-edited XML). Checking
# both means neither a localized Word install nor a non-Word export loses
# chapter structure; a document matching neither still falls back to
# reading as a single "Full Document" chapter (the same degraded-but-not-
# broken outcome as a PDF with no outline) -- not treated as a parse
# failure.
_DOCX_HEADING_STYLE_NAMES = {"Title", "Heading 1", "Heading 2", "Heading 3"}


def _table_text(table) -> str:
    """Flattens a table to plain text: one line per row, cells tab-separated.

    Shared by both `_parse_docx` and `_parse_pptx` -- a DOCX `Table` and a
    PPTX `GraphicFrame.table` are unrelated classes from unrelated
    libraries, but both expose the same `.rows` of cells with a `.text`
    property, so one function handles either.

    A row whose cells are all empty (or a table with none) contributes
    nothing -- same "skip empty" treatment as a blank paragraph gets
    below, so an empty table doesn't inject blank lines into a chapter's
    text.

    Deduped by `id(cell)` before reading `.text`: a horizontally-merged
    DOCX cell comes back as the *same* `_Cell` instance at every grid
    column it spans (word represents the merge by grid-column-count, not
    by a separate cell per column), so reading every column verbatim
    repeated the merged text once per spanned column ("Merged header
    Merged header Merged header"). Deduped by identity rather than
    `dict.fromkeys(row.cells)` because python-pptx's `_Cell` has no
    `__hash__` of its own (raises `TypeError: unhashable type`) --
    `id()` needs no cooperation from either library's cell class. PPTX
    merges work differently anyway -- a spanned cell's own `.text` is
    already empty -- so this dedup is a harmless no-op there, not a
    second code path.
    """
    lines = []
    for row in table.rows:
        seen_cell_ids: set[int] = set()
        cell_texts = []
        for cell in row.cells:
            if id(cell) in seen_cell_ids:
                continue
            seen_cell_ids.add(id(cell))
            text = cell.text.strip()
            if text:
                cell_texts.append(text)
        line = "\t".join(cell_texts)
        if line:
            lines.append(line)
    return "\n".join(lines)


def _parse_docx(content: bytes) -> list[tuple[str, str]]:
    _check_zip_bomb(content)
    try:
        document = DocxDocument(io.BytesIO(content))
    except Exception as exc:
        raise UnparseableDocument("Could not parse DOCX.") from exc

    chapters: list[tuple[str, list[str]]] = []
    current_title = _FULL_DOCUMENT_CHAPTER
    current_parts: list[str] = []

    # `iter_inner_content()` (not `document.paragraphs`) walks paragraphs
    # *and* top-level tables together in actual document order -- a table
    # is exactly as common in a Word document as a paragraph, so reading
    # only paragraphs silently dropped every table's text, and a document
    # whose content lived entirely in a table used to fail to parse at
    # all (zero extractable text).
    for block in document.iter_inner_content():
        if isinstance(block, DocxTable):
            table_text = _table_text(block)
            if table_text:
                current_parts.append(table_text)
            continue

        text = block.text.strip()
        if not text:
            continue
        style = block.style
        is_heading = style is not None and (
            style.style_id in _DOCX_HEADING_STYLE_IDS
            or style.name in _DOCX_HEADING_STYLE_NAMES
        )
        if is_heading:
            # Appended even when `current_parts` is empty -- a heading
            # directly followed by another heading (no body paragraph
            # between them, e.g. a "Part I" divider before "Chapter 1")
            # otherwise vanished here before it ever reached
            # `_chunk_chapters`'s own empty-body handling. The initial
            # placeholder title is the one exception: an empty preamble
            # before the first real heading isn't a heading a user wrote,
            # so there's nothing worth keeping a chapter for.
            if current_parts or current_title != _FULL_DOCUMENT_CHAPTER:
                chapters.append((current_title, current_parts))
            current_title = text
            current_parts = []
        else:
            current_parts.append(text)

    if current_parts or current_title != _FULL_DOCUMENT_CHAPTER:
        chapters.append((current_title, current_parts))

    return [(title, "\n".join(parts)) for title, parts in chapters]


# --- PPTX -----------------------------------------------------------------


def _parse_pptx(content: bytes) -> list[tuple[str, str]]:
    _check_zip_bomb(content)
    try:
        presentation = Presentation(io.BytesIO(content))
    except Exception as exc:
        raise UnparseableDocument("Could not parse PPTX.") from exc

    chapters: list[tuple[str, str]] = []
    for index, slide in enumerate(presentation.slides, start=1):
        title_shape = slide.shapes.title
        title_text = (
            title_shape.text_frame.text.strip()
            if title_shape is not None and title_shape.has_text_frame
            else ""
        )

        parts: list[str] = []
        for shape in slide.shapes:
            # Identified by placeholder idx, not `shape is title_shape`:
            # `slide.shapes` builds a fresh wrapper object on every
            # iteration, so an identity check against `title_shape` (built
            # by a separate `.title` lookup above) would never match even
            # for the same underlying slide element.
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                continue
            # Checked before `has_text_frame`: a table shape is a
            # `GraphicFrame`, which reads as `has_text_frame == False` (it
            # holds a table, not a text frame), so table slides used to
            # contribute nothing at all -- the same "tables are as common
            # as body text" gap `_parse_docx` had.
            if shape.has_table:
                table_text = _table_text(shape.table)
                if table_text:
                    parts.append(table_text)
                continue
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if text:
                parts.append(text)

        # Appended unconditionally, even when `parts` is empty: a
        # title-only slide (a section divider, an agenda slide) is still
        # a real, findable slide -- `_chunk_chapters` falls back to
        # indexing the title itself when a chapter has no body text.
        #
        # `f"Slide {index}"` is a *generated* label, not something a user
        # wrote -- it must never reach `_chunk_chapters`'s title-as-text
        # fallback the way a real title does, or an image-only slide (no
        # title, no text) ends up "indexed" as the literal string "Slide
        # 1" with nothing behind it. So it's only used when there's real
        # body text to label; with no title and no body, the chapter
        # falls back to `_FULL_DOCUMENT_CHAPTER` instead -- the same
        # sentinel a genuinely empty PDF page produces, which
        # `_chunk_chapters` already knows to drop rather than index.
        if title_text:
            chapter_title = title_text
        elif parts:
            chapter_title = f"Slide {index}"
        else:
            chapter_title = _FULL_DOCUMENT_CHAPTER
        chapters.append((chapter_title, "\n".join(parts)))

    return chapters


# --- PDF ------------------------------------------------------------------


# Fraction of whitespace-separated tokens that must be a single character
# before a page is treated as glyph-spaced (see `_repair_char_spacing`).
# High deliberately: real prose in any language sits far below this (even
# Chinese/Japanese text extracts as multi-character tokens here), so a
# genuine document cannot trip it by accident. The observed failing PDFs
# score 1.00 -- every single token -- so there is no need to sit close to
# the boundary.
_CHAR_SPACED_TOKEN_RATIO = 0.8

# Below this many tokens the ratio isn't meaningful -- a near-empty page
# ("A", "B") would otherwise look 100% char-spaced and get "repaired"
# into nonsense.
_CHAR_SPACED_MIN_TOKENS = 20


def _looks_char_spaced(text: str) -> bool:
    """True if `text` looks like every glyph was extracted separately.

    Some PDFs (design-tool exports especially) position each glyph
    individually rather than as runs of text. pypdf faithfully reports
    what the file says, so "Контакти" comes back as "К о н т а к т и".
    """
    tokens = text.split()
    if len(tokens) < _CHAR_SPACED_MIN_TOKENS:
        return False
    singles = sum(1 for token in tokens if len(token) == 1)
    return singles / len(tokens) >= _CHAR_SPACED_TOKEN_RATIO


def _repair_char_spacing(text: str) -> str:
    """Rejoins glyph-spaced text into words.

    Relies on the one distinction such PDFs do preserve: a single space
    separates glyphs *within* a word, while a real word boundary comes
    through as two or more spaces. So "г р а д  К ю с т е н д и л"
    carries enough information to recover "град Кюстендил" exactly --
    this is a lossless reconstruction, not a guess at where words go.
    Line structure is preserved so chapter-heading detection downstream
    still sees the same lines.

    Only ever called behind `_looks_char_spaced`, because applied to
    normal text it would destroy every genuine single-letter word ("a",
    "I") by gluing it onto its neighbour.
    """
    repaired_lines = []
    for line in text.splitlines():
        if not line.strip():
            repaired_lines.append(line)
            continue
        words = [re.sub(r"\s+", "", word) for word in re.split(r"\s{2,}", line.strip())]
        repaired_lines.append(" ".join(word for word in words if word))
    return "\n".join(repaired_lines)


def _parse_pdf(content: bytes) -> list[tuple[str, str]]:
    try:
        reader = PdfReader(io.BytesIO(content))
        # Per page, not over the whole document: a PDF can mix a
        # glyph-spaced cover/graphic page with normally-extracted body
        # pages, and a document-wide ratio would let the good pages mask
        # a bad one (or vice versa).
        page_texts = []
        for page in reader.pages:
            page_text = page.extract_text() or ""
            if _looks_char_spaced(page_text):
                page_text = _repair_char_spacing(page_text)
            page_texts.append(page_text)
    except Exception as exc:
        raise UnparseableDocument("Could not parse PDF.") from exc

    chapters = _pdf_chapters_from_outline(reader, page_texts)
    if chapters is not None:
        return chapters
    return [(_FULL_DOCUMENT_CHAPTER, "\n".join(page_texts))]


def _pdf_chapters_from_outline(
    reader: PdfReader, page_texts: list[str]
) -> list[tuple[str, str]] | None:
    """Top-level PDF bookmarks as chapters, or `None` to fall back to a
    single "Full Document" chapter.

    A strict, checkable condition rather than an open-ended "if it looks
    clean" judgment call: the outline must be a flat list (no nested
    sub-lists) and every entry must resolve to a page number with no
    exception. Any nesting, an empty outline, or any resolution failure
    falls back immediately -- no partial/best-effort read of a broken or
    deeply-nested table of contents.
    """
    try:
        outline = reader.outline
    except Exception:
        return None
    if not outline:
        return None
    if any(isinstance(item, list) for item in outline):
        return None

    entries: list[tuple[str, int]] = []
    try:
        for item in outline:
            page_number = reader.get_destination_page_number(item)
            if page_number is None:
                return None
            title = item.title or _FULL_DOCUMENT_CHAPTER
            entries.append((title, page_number))
    except Exception:
        return None

    entries.sort(key=lambda entry: entry[1])
    chapters: list[tuple[str, str]] = []
    if entries[0][1] > 0:
        # Pages before the first bookmark (title page, table of contents,
        # an abstract) are real content the outline just doesn't label --
        # not something to drop because no chapter claims them.
        preamble = "\n".join(page_texts[: entries[0][1]])
        if preamble.strip():
            chapters.append((_FULL_DOCUMENT_CHAPTER, preamble))
    for i, (title, start_page) in enumerate(entries):
        end_page = entries[i + 1][1] if i + 1 < len(entries) else len(page_texts)
        chapters.append((title, "\n".join(page_texts[start_page:end_page])))
    return chapters
