"""Unit tests for `app.documents.parsing` -- format-specific parsing bugs
that don't need the full upload/ingest flow to reproduce or pin down.
"""

import io
import zipfile
import zlib

import pytest
from docx import Document as DocxDocument
from pptx import Presentation

from app.documents import parsing
from app.documents.parsing import UnparseableDocument, parse_document


def _docx_bytes(build) -> bytes:
    document = DocxDocument()
    build(document)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _pptx_bytes(build) -> bytes:
    presentation = Presentation()
    build(presentation)
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()

_MINIMAL_PDF_WITH_PREAMBLE = b"""%PDF-1.4
1 0 obj<< /Type /Catalog /Pages 2 0 R /Outlines 6 0 R >>endobj
2 0 obj<< /Type /Pages /Kids [3 0 R 7 0 R] /Count 2 >>endobj
3 0 obj<< /Type /Page /Parent 2 0 R /Resources << /Font << /F1 4 0 R >> >> /MediaBox [0 0 200 200] /Contents 5 0 R >>endobj
4 0 obj<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>endobj
5 0 obj<< /Length 50 >>
stream
BT /F1 12 Tf 10 100 Td (Title page preamble text) Tj ET
endstream
endobj
6 0 obj<< /Type /Outlines /First 8 0 R /Last 8 0 R /Count 1 >>endobj
7 0 obj<< /Type /Page /Parent 2 0 R /Resources << /Font << /F1 4 0 R >> >> /MediaBox [0 0 200 200] /Contents 9 0 R >>endobj
8 0 obj<< /Title (Chapter One) /Parent 6 0 R /Dest [7 0 R /Fit] >>endobj
9 0 obj<< /Length 40 >>
stream
BT /F1 12 Tf 10 100 Td (Chapter one body) Tj ET
endstream
endobj
xref
0 10
trailer<< /Size 10 /Root 1 0 R >>
startxref
0
%%EOF"""


def test_markdown_preserves_text_before_first_heading():
    chunks = parse_document(
        "markdown", b"Intro paragraph that matters.\n\n# Chapter One\n\nBody here."
    )

    assert any("Intro paragraph that matters." in c.text for c in chunks)
    preamble_chunk = next(c for c in chunks if "Intro paragraph that matters." in c.text)
    assert preamble_chunk.chapter == "Full Document"
    assert preamble_chunk.chunk_index == 0

    chapter_one_chunk = next(c for c in chunks if c.chapter == "Chapter One")
    assert "Body here." in chapter_one_chunk.text


def test_markdown_with_no_preamble_has_no_spurious_full_document_chapter():
    chunks = parse_document("markdown", b"# Chapter One\n\nBody here.")

    assert {c.chapter for c in chunks} == {"Chapter One"}


def test_pdf_outline_preserves_pages_before_first_bookmark():
    chunks = parse_document("pdf", _MINIMAL_PDF_WITH_PREAMBLE)

    preamble_chunk = next(c for c in chunks if "Title page preamble text" in c.text)
    assert preamble_chunk.chapter == "Full Document"
    assert preamble_chunk.chunk_index == 0

    chapter_chunk = next(c for c in chunks if c.chapter == "Chapter One")
    assert "Chapter one body" in chapter_chunk.text
    assert chapter_chunk.chunk_index == 1


def test_html_comments_are_excluded_from_passage_text():
    chunks = parse_document(
        "html",
        b"<html><body><!-- secret internal note --><h1>T</h1><p>Body</p></body></html>",
    )

    assert len(chunks) == 1
    assert chunks[0].chapter == "T"
    assert chunks[0].text == "Body"
    assert "secret internal note" not in chunks[0].text


def test_html_comment_before_first_heading_does_not_leak_into_preamble():
    chunks = parse_document(
        "html",
        b"<html><body><!-- internal note -->Real preamble text<h1>T</h1><p>Body</p></body></html>",
    )

    all_text = " ".join(c.text for c in chunks)
    assert "internal note" not in all_text
    assert "Real preamble text" in all_text


def test_chunk_word_count_stays_within_the_embedding_model_s_token_budget():
    # Sized against shared/embeddings/model.py's 512-token multilingual
    # limit -- not re-asserting the exact constant, but that whatever it
    # is stays comfortably under a budget that would otherwise mean most
    # of a chunk's stored, citable text was never actually embedded.
    text = " ".join(f"word{i}" for i in range(1000))
    chunks = parse_document("markdown", text.encode())

    for chunk in chunks:
        assert len(chunk.text.split()) <= parsing._CHUNK_WORD_COUNT


def test_markdown_heading_inside_a_code_fence_is_not_treated_as_a_chapter():
    md = (
        b"# Real Heading\n\n"
        b"Some intro text.\n\n"
        b"```bash\n# install deps\npip install foo\n```\n\n"
        b"## Another Real Heading\n\nFinal text."
    )

    chunks = parse_document("markdown", md)

    chapters = {c.chapter for c in chunks}
    assert chapters == {"Real Heading", "Another Real Heading"}
    assert "install deps" not in chapters


# ---------------------------------------------------------------------------
# DOCX / PPTX.
# ---------------------------------------------------------------------------


def test_docx_headings_become_chapters():
    def build(document):
        document.add_paragraph("Intro before any heading.")
        document.add_heading("Chapter One", level=1)
        document.add_paragraph("Body of chapter one.")
        document.add_heading("Chapter Two", level=2)
        document.add_paragraph("Body of chapter two.")

    chunks = parse_document("docx", _docx_bytes(build))

    assert {c.chapter for c in chunks} == {"Full Document", "Chapter One", "Chapter Two"}
    chapter_one_chunk = next(c for c in chunks if c.chapter == "Chapter One")
    assert "Body of chapter one." in chapter_one_chunk.text
    chapter_two_chunk = next(c for c in chunks if c.chapter == "Chapter Two")
    assert "Body of chapter two." in chapter_two_chunk.text


def test_docx_headings_become_chapters_with_localized_style_names():
    # A non-English-localized Word (e.g. German) renames the *displayed*
    # style name (here simulated directly, since python-docx itself always
    # writes English names) while keeping the built-in style_id
    # ("Heading1") unchanged -- headings must still be detected via that
    # locale-invariant id, not the localized display name.
    def build(document):
        document.add_paragraph("Intro before any heading.")
        document.add_heading("Kapitel Eins", level=1)
        document.add_paragraph("Text von Kapitel eins.")
        document.styles["Heading 1"].name = "Überschrift 1"

    chunks = parse_document("docx", _docx_bytes(build))

    assert {c.chapter for c in chunks} == {"Full Document", "Kapitel Eins"}
    chapter_chunk = next(c for c in chunks if c.chapter == "Kapitel Eins")
    assert "Text von Kapitel eins." in chapter_chunk.text


def test_docx_with_no_headings_has_single_full_document_chapter():
    def build(document):
        document.add_paragraph("Just a plain paragraph.")

    chunks = parse_document("docx", _docx_bytes(build))

    assert {c.chapter for c in chunks} == {"Full Document"}


def test_docx_unparseable_bytes_raise():
    with pytest.raises(UnparseableDocument):
        parse_document("docx", b"not a real docx file")


def test_pptx_slide_title_becomes_chapter_and_is_not_duplicated_into_body():
    def build(presentation):
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        slide.shapes.title.text = "Welcome Slide"
        slide.placeholders[1].text = "Subtitle body text."

    chunks = parse_document("pptx", _pptx_bytes(build))

    assert {c.chapter for c in chunks} == {"Welcome Slide"}
    assert any("Subtitle body text." in c.text for c in chunks)
    assert not any("Welcome Slide" in c.text for c in chunks)


def test_pptx_slide_without_title_gets_generated_chapter_name():
    def build(presentation):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        textbox = slide.shapes.add_textbox(0, 0, 100, 100)
        textbox.text_frame.text = "Body text with no title placeholder."

    chunks = parse_document("pptx", _pptx_bytes(build))

    assert {c.chapter for c in chunks} == {"Slide 1"}
    assert any("Body text with no title placeholder." in c.text for c in chunks)


def test_pptx_unparseable_bytes_raise():
    with pytest.raises(UnparseableDocument):
        parse_document("pptx", b"not a real pptx file")


# ---------------------------------------------------------------------------
# DOCX / PPTX tables. Previously unread entirely (`document.paragraphs`
# never visits a table, and a table shape reads as `has_text_frame ==
# False`), so a document whose content lived only in a table produced zero
# extractable text -- a hard `UnparseableDocument` failure, not a
# degradation.
# ---------------------------------------------------------------------------


def test_docx_table_text_is_extracted():
    def build(document):
        document.add_paragraph("Intro paragraph.")
        table = document.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Name"
        table.cell(0, 1).text = "Role"
        table.cell(1, 0).text = "Alice"
        table.cell(1, 1).text = "Engineer"

    chunks = parse_document("docx", _docx_bytes(build))

    all_text = " ".join(c.text for c in chunks)
    assert "Alice" in all_text
    assert "Engineer" in all_text


def test_docx_table_only_document_does_not_fail_to_parse():
    def build(document):
        table = document.add_table(rows=1, cols=2)
        table.cell(0, 0).text = "Name"
        table.cell(0, 1).text = "Bob"

    chunks = parse_document("docx", _docx_bytes(build))

    assert any("Bob" in c.text for c in chunks)


def test_docx_merged_cell_text_is_not_repeated_per_spanned_column():
    # `row.cells` hands back the *same* `_Cell` object at every grid
    # column a horizontal merge spans (word represents the merge by
    # grid-column-count, not by one cell per column) -- reading each
    # column verbatim used to repeat the merged text once per spanned
    # column ("Merged header Merged header Merged header a b c").
    def build(document):
        table = document.add_table(rows=2, cols=3)
        table.cell(0, 0).merge(table.cell(0, 2))
        table.cell(0, 0).text = "Merged header"
        table.cell(1, 0).text = "a"
        table.cell(1, 1).text = "b"
        table.cell(1, 2).text = "c"

    chunks = parse_document("docx", _docx_bytes(build))

    all_text = " ".join(c.text for c in chunks)
    assert all_text.count("Merged header") == 1


def test_docx_heading_with_no_body_paragraph_is_not_dropped():
    # A heading directly followed by another heading (a "Part I" divider
    # before "Chapter 1") used to vanish: the chapter-switch append was
    # gated on `current_parts` being non-empty, so "Part I" was discarded
    # before it ever became a chunk.
    def build(document):
        document.add_heading("Part I", level=1)
        document.add_heading("Chapter 1", level=1)
        document.add_paragraph("Body.")

    chunks = parse_document("docx", _docx_bytes(build))

    assert {c.chapter for c in chunks} == {"Part I", "Chapter 1"}
    part_one_chunk = next(c for c in chunks if c.chapter == "Part I")
    assert "Part I" in part_one_chunk.text


def test_pptx_table_text_is_extracted():
    def build(presentation):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        graphic_frame = slide.shapes.add_table(2, 2, 0, 0, 100, 100)
        table = graphic_frame.table
        table.cell(0, 0).text = "Name"
        table.cell(0, 1).text = "Role"
        table.cell(1, 0).text = "Carol"
        table.cell(1, 1).text = "Designer"

    chunks = parse_document("pptx", _pptx_bytes(build))

    all_text = " ".join(c.text for c in chunks)
    assert "Carol" in all_text
    assert "Designer" in all_text


def test_pptx_title_only_slide_is_still_indexed():
    # A section-divider/agenda slide with only a title used to disappear:
    # its chapter's body text is empty, and `_chunk_chapters` dropped any
    # chapter with no words -- so a search for "Next Steps" could never
    # find the slide that says exactly that.
    def build(presentation):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "Next Steps"

    chunks = parse_document("pptx", _pptx_bytes(build))

    assert any(c.chapter == "Next Steps" for c in chunks)
    assert any("Next Steps" in c.text for c in chunks)


def test_pptx_image_only_slide_fails_to_parse_rather_than_indexing_its_placeholder_label():
    # Regression: `f"Slide {index}"` is a *generated* label, not
    # something a user wrote. It used to be handed to `_chunk_chapters`
    # exactly like a real title, so an image-only slide (no title
    # placeholder, no text anywhere on it) ended up "indexed" as the
    # literal, contentless string "Slide 1" -- retrievable and citable in
    # chat with nothing real behind it. `parse_document`'s own contract
    # (mirroring a scanned/image-only PDF) is that a document with zero
    # real extractable text must fail, not silently succeed with a
    # placeholder standing in for content.
    def build(presentation):
        presentation.slides.add_slide(presentation.slide_layouts[6])
        presentation.slides.add_slide(presentation.slide_layouts[6])

    with pytest.raises(UnparseableDocument):
        parse_document("pptx", _pptx_bytes(build))


def test_pptx_slide_with_body_text_but_no_title_still_gets_a_slide_number():
    # Companion to the image-only case above: a slide with real body text
    # but no title placeholder must still get its `f"Slide {index}"`
    # label -- only a slide with *no* real content anywhere falls back to
    # the "nothing here" sentinel.
    def build(presentation):
        presentation.slides.add_slide(presentation.slide_layouts[6])  # no title, no body
        slide2 = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide2.shapes.add_textbox(0, 0, 100, 100).text_frame.text = "Real body text."

    chunks = parse_document("pptx", _pptx_bytes(build))

    assert {c.chapter for c in chunks} == {"Slide 2"}
    assert any("Real body text." in c.text for c in chunks)


# ---------------------------------------------------------------------------
# DOCX / PPTX zip-bomb guard. Unlike every other supported format,
# MAX_FILE_SIZE_BYTES only bounds the *compressed* upload -- DEFLATE lets a
# small archive expand enormously in memory before python-docx/python-pptx
# get a chance to fail gracefully.
#
# `_MAX_ZIP_UNCOMPRESSED_BYTES` is monkeypatched down to a tiny cap in
# every test below rather than actually building a payload past the real
# 200MB one -- the mechanism being tested (does the guard measure real
# decompressed bytes and stop the moment the cap is crossed) doesn't
# depend on the cap's real-world size, and building/holding an
# honest-to-goodness 200MB+ payload in the test process for every run
# would be exactly the kind of memory cost this guard exists to avoid.
# ---------------------------------------------------------------------------


def _zip_bomb_bytes(*, size: int) -> bytes:
    """A single-entry zip archive whose one file inflates to `size` bytes
    of a single repeated byte -- trivially compressible, so `size` can be
    made arbitrarily larger than the archive's own compressed byte count."""
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("word/document.xml", b"A" * size)
    return buffer.getvalue()


def _lie_about_declared_uncompressed_size(archive_bytes: bytes, *, declared: int) -> bytes:
    """Rewrites a single-entry archive's declared uncompressed size, in
    both the local file header and the central directory record, to
    `declared` -- without touching the real compressed data at all.

    Simulates the attack `_check_zip_bomb` has to survive: `ZipInfo
    .file_size` is read straight from the central directory, which is
    ordinary attacker-controlled archive data, not a measurement `zipfile`
    verifies against the actual compressed stream until (if ever) a
    caller reads the whole entry and its CRC-32 is checked. A crafted
    archive can carry a real, large deflate stream behind a central
    directory entry that simply declares a small size.
    """
    data = bytearray(archive_bytes)
    assert data[0:4] == b"PK\x03\x04", "expected a local file header at offset 0"
    data[22:26] = declared.to_bytes(4, "little")  # local header uncompressed-size field
    # `rfind`, not `find`: the *real* central directory record is always
    # the last occurrence of this signature, immediately before the
    # end-of-central-directory record -- a highly compressible payload
    # like this one's repeated byte can coincidentally contain the same
    # 4-byte signature earlier, inside the compressed data itself, and
    # patching that false match would corrupt the archive instead.
    central_directory_offset = archive_bytes.rfind(b"PK\x01\x02")
    assert central_directory_offset != -1, "expected a central directory record"
    data[central_directory_offset + 24 : central_directory_offset + 28] = declared.to_bytes(
        4, "little"
    )  # central directory record's own uncompressed-size field
    return bytes(data)


def test_docx_zip_bomb_is_rejected_before_expanding_in_memory(monkeypatch):
    monkeypatch.setattr(parsing, "_MAX_ZIP_UNCOMPRESSED_BYTES", 1_000)

    with pytest.raises(UnparseableDocument, match="unreasonable size"):
        parse_document("docx", _zip_bomb_bytes(size=5_000))


def test_pptx_zip_bomb_is_rejected_before_expanding_in_memory(monkeypatch):
    monkeypatch.setattr(parsing, "_MAX_ZIP_UNCOMPRESSED_BYTES", 1_000)

    with pytest.raises(UnparseableDocument, match="unreasonable size"):
        parse_document("pptx", _zip_bomb_bytes(size=5_000))


def test_zip_bomb_guard_rejects_rather_than_trusts_a_lying_declared_size(monkeypatch):
    # The regression this guards against: the guard's first version
    # summed `ZipInfo.file_size` -- trusting the archive's own metadata,
    # which is attacker-controlled -- rather than reading any actual
    # bytes. This archive's central directory *under*-declares its real
    # decompressed size (claims 10 bytes; the real payload is 5000,
    # unchanged), which is exactly what that first version would have
    # been fooled by.
    #
    # Streaming the entry is necessary but was not sufficient on its own:
    # `zipfile` applies the declared size to the *read* too
    # (`ZipExtFile` truncates each read to the byte count it has left),
    # so a guard that reads the entry normally measures the lie -- 10
    # bytes, comfortably under the cap -- and waves the archive through.
    # Reading against a `ZipInfo` copy whose declared size can't truncate
    # anything is what makes the count reflect the real 5000, so the size
    # error fires here rather than the archive being accepted as "10
    # bytes, safely under the cap".
    monkeypatch.setattr(parsing, "_MAX_ZIP_UNCOMPRESSED_BYTES", 1_000)
    honest_archive = _zip_bomb_bytes(size=5_000)
    lying_archive = _lie_about_declared_uncompressed_size(honest_archive, declared=10)

    with pytest.raises(UnparseableDocument, match="unreasonable size"):
        parse_document("docx", lying_archive)


def test_zip_bomb_guard_rejects_a_lying_size_whose_crc_was_forged_to_match_it(monkeypatch):
    # The one way to keep the archive internally consistent while lying
    # about the size: forge the CRC-32 to match only the *declared*
    # (small) prefix of the real data, so a reader that stops at the
    # declared size sees a checksum that agrees with what it got.
    #
    # That lie is worth defending against because the truncation it
    # exploits bounds what a reader *receives*, not what inflating
    # allocated to produce it -- `ZipFile.read(name)`, which is what
    # python-docx/python-pptx call, passes no size and so decompresses
    # with an effectively unbounded `max_length`, materializing the whole
    # payload before slicing it back down. A guard fooled here would hand
    # those libraries an archive that peaks at hundreds of MB of heap
    # (measured: 830MB from a 433KB upload) on its way to failing.
    #
    # Reading past the declared size is what exposes it: the forged CRC
    # only matches the prefix, so a full read of the real payload no
    # longer agrees with it and `zipfile` itself rejects the entry -- the
    # forgery that hides the size is what makes the checksum fail. Which
    # of the two rejections lands depends only on scale: this payload is
    # small enough that one chunk read reaches the end of the entry (and
    # so its CRC check) first, while a real bomb crosses the size cap
    # chunks before the decompressor ever reaches eof.
    monkeypatch.setattr(parsing, "_MAX_ZIP_UNCOMPRESSED_BYTES", 1_000)
    honest_archive = bytearray(_zip_bomb_bytes(size=5_000))
    declared = 10
    forged_crc = zlib.crc32(b"A" * declared)
    honest_archive[14:18] = forged_crc.to_bytes(4, "little")  # local header CRC-32
    central_directory_offset = bytes(honest_archive).rfind(b"PK\x01\x02")
    honest_archive[central_directory_offset + 16 : central_directory_offset + 20] = (
        forged_crc.to_bytes(4, "little")
    )  # central directory record's own CRC-32
    lying_archive = _lie_about_declared_uncompressed_size(bytes(honest_archive), declared=declared)

    with pytest.raises(UnparseableDocument, match="Could not parse archive"):
        parse_document("docx", lying_archive)


def test_zip_bomb_guard_leaves_the_archives_own_entry_metadata_untouched():
    # The guard measures against a *copy* of each `ZipInfo` -- the ones
    # `infolist()` returns are the live objects the `ZipFile` uses for
    # its own reads, so overwriting `file_size` on them to defeat
    # truncation would corrupt the archive's view of every entry for
    # readers that come after (python-docx, immediately after this).
    archive_bytes = _zip_bomb_bytes(size=5_000)

    parsing._check_zip_bomb(archive_bytes)

    with zipfile.ZipFile(io.BytesIO(archive_bytes)) as archive:
        assert [info.file_size for info in archive.infolist()] == [5_000]
        assert archive.read("word/document.xml") == b"A" * 5_000


def test_legitimate_docx_well_under_the_zip_bomb_cap_still_parses():
    def build(document):
        document.add_paragraph("An ordinary paragraph.")

    chunks = parse_document("docx", _docx_bytes(build))

    assert any("An ordinary paragraph." in c.text for c in chunks)


# ---------------------------------------------------------------------------
# Glyph-spaced PDF repair. Some PDFs (design-tool exports especially)
# position every character individually, so pypdf faithfully returns
# "К о н т а к т и" for "Контакти". Observed on two real CV PDFs where
# 100% of extracted tokens were single characters.
# ---------------------------------------------------------------------------

from app.documents.parsing import _looks_char_spaced, _repair_char_spacing


def test_char_spacing_detected_on_fully_glyph_spaced_text():
    text = "К о н т а к т и\n0 8 8 6 9 9 7 8 5 8\nТ е л е ф о н\ny o a n a s b @ g m a i l . c o m"
    assert _looks_char_spaced(text)


def test_normal_prose_is_not_detected_as_char_spaced():
    """The detector must not fire on real documents -- repairing one would
    glue every word to its neighbour."""
    text = (
        "Project Aurora is GraphMind's internal codename for the Q3 2026 "
        "knowledge-graph visualization redesign. The project began on "
        "2026-04-01 and is scheduled to ship on 2026-09-15. Elena Rusev "
        "leads it, and Northwind Robotics supplies the hardware."
    )
    assert not _looks_char_spaced(text)


def test_bulgarian_prose_is_not_detected_as_char_spaced():
    """Cyrillic is the script the failing PDFs were in -- the detector must
    key on token length, not on the alphabet."""
    text = (
        "Личен дневник по разработката на Gamification модула върху "
        "стажантската среда на Sirma Academy LMS. Всеки запис отговаря на "
        "commit — какво е добавено и защо, подредено хронологично."
    )
    assert not _looks_char_spaced(text)


def test_short_text_is_never_treated_as_char_spaced():
    """Below the token floor the ratio is meaningless -- a two-token page
    would otherwise score 100% and be 'repaired' into nonsense."""
    assert not _looks_char_spaced("A B")
    assert not _looks_char_spaced("")


def test_repair_rejoins_words_using_the_double_space_boundary():
    """Single space separates glyphs; two or more marks a real word
    boundary. That distinction is what makes this reconstruction lossless
    rather than a guess."""
    assert _repair_char_spacing("г р а д  К ю с т е н д и л") == "град Кюстендил"
    assert _repair_char_spacing("Y O A N A  B O R I S O V A") == "YOANA BORISOVA"


def test_repair_preserves_line_structure():
    """Chapter-heading detection downstream still reads these lines, so
    the line count must not change."""
    repaired = _repair_char_spacing("К о н т а к т и\n0 8 8 6 9 9 7 8 5 8")
    assert repaired == "Контакти\n0886997858"


def test_repair_keeps_punctuation_inside_a_rejoined_token():
    assert _repair_char_spacing("y o a n a s b @ g m a i l . c o m") == "yoanasb@gmail.com"


def test_repaired_pdf_text_is_substantially_shorter():
    """The repair roughly halves the character count, which is why it also
    matters for cost: extraction text is capped at EXTRACTION_CHAR_BUDGET,
    so glyph spacing was spending half that budget on spaces."""
    spaced = "\n".join(["К о н т а к т и  Т е л е ф о н"] * 10)
    assert len(_repair_char_spacing(spaced)) < len(spaced) / 1.8
